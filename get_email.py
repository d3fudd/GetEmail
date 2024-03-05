import os
import re
import docx
import pptx
import openpyxl
import fitz
import magic

def main():
    directory = input("Digite o caminho do diretório: ")
    emails = find_emails_in_directory(directory)
    unique_emails = set(emails)
    print("\n================================================\n")
    print("E-mails encontrados:")
    for email in unique_emails:
        print(email)

def find_emails_in_directory(directory):
    emails_found = []
    for root, _, files in os.walk(directory):
        for file in files:
            file_path = os.path.join(root, file)
            file_type = get_file_type(file_path)
            if file_type == 'docx':
                try:
                    emails_found.extend(extract_emails_from_docx(file_path))
                except Exception as e:
                    print(f"Erro ao processar o arquivo {file_path}: {e}")

            elif file_type == 'pptx':
                try:
                    emails_found.extend(extract_emails_from_pptx(file_path))
                except Exception as e:
                    print(f"Erro ao processar o arquivo {file_path}: {e}")

            elif file_type == 'xlsx':
                try:
                    emails_found.extend(extract_emails_from_xlsx(file_path))
                except Exception as e:
                    print(f"Erro ao processar o arquivo {file_path}: {e}")

            elif file_type == 'pdf':
                try:
                    emails_found.extend(extract_emails_from_pdf(file_path))
                except Exception as e:
                    print(f"Erro ao processar o arquivo {file_path}: {e}")

            elif file_type == 'odt':
                try:
                    emails_found.extend(extract_emails_from_odt(file_path))
                except Exception as e:
                    print(f"Erro ao processar o arquivo {file_path}: {e}")

            elif file_type == 'doc':
                try:
                    emails_found.extend(extract_emails_from_doc(file_path))
                except Exception as e:
                    print(f"Erro ao processar o arquivo {file_path}: {e}")

            else:
                print(f"Formato não suportado para o arquivo {file_path}")
    return emails_found

def get_file_type(file_path):
    mime = magic.Magic(mime=True)
    file_mime = mime.from_file(file_path)
    if 'docx' in file_mime:
        return 'docx'
    elif 'pptx' in file_mime:
        return 'pptx'
    elif 'xlsx' in file_mime:
        return 'xlsx'
    elif 'pdf' in file_mime:
        return 'pdf'
    elif 'odt' in file_mime:
        return 'odt'
    elif 'doc' in file_mime:
        return 'doc'
    else:
        return None

def extract_emails_from_docx(file_path):
    emails = []
    try:
        doc = docx.Document(file_path)
        for paragraph in doc.paragraphs:
            emails.extend(re.findall(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', paragraph.text))
    except Exception as e:
        print(f"Erro ao processar o arquivo {file_path}: {e}")
    return emails

def extract_emails_from_doc(file_path):
    emails = []
    try:
        doc = doc.Document(file_path)
        for paragraph in doc.paragraphs:
            emails.extend(re.findall(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', paragraph.text))
    except Exception as e:
        print(f"Erro ao processar o arquivo {file_path}: {e}")
    return emails

def extract_emails_from_pptx(file_path):
    emails = []
    presentation = pptx.Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                emails.extend(re.findall(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', shape.text))
    return emails

def extract_emails_from_xlsx(file_path):
    emails = []
    wb = openpyxl.load_workbook(file_path)
    for sheet in wb:
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    emails.extend(re.findall(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', cell.value))
    return emails

def extract_emails_from_odt(file_path):
    emails = []
    with open(file_path, 'rb') as file:
        doc = docx.opendocx(file)
        for paragraph in doc:
            emails.extend(re.findall(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', paragraph.text))
    return emails

def extract_emails_from_pdf(file_path):
    emails = []
    with fitz.open(file_path) as doc:
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            emails.extend(re.findall(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', text))
    return emails

if __name__ == "__main__":
    main()
